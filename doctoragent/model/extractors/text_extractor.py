"""Text and document content extractor.

Supports the following sources:

- Plain text (``txt/md/csv/json/log/xml/html``): read directly as UTF-8,
  falling back to ``latin-1``.
- ``pdf``: extracts text via ``pypdf`` (a lightweight pure-Python library),
  with optional advanced table/layout/metadata extraction via ``pdfplumber``.
- ``docx``: extracts paragraph text and structure via ``python-docx``.
- ``xlsx``: extracts cell text and tables via ``openpyxl``.
- ``pptx``: extracts slide text via ``python-pptx``.
- ``eml``: extracts email headers and body via the stdlib ``email`` module.
- ``rtf``: extracts plain text via ``striprtf``.

Third-party libraries (``pypdf``, ``python-docx``, ``pdfplumber``,
``openpyxl``, ``python-pptx``, ``striprtf``) belong to the ``[multimodal]``
extras. When a dependency is missing, ``supports`` still returns ``True`` by
extension, and ``extract`` returns empty text + ``method="none"`` without
raising, so the caller can degrade gracefully.

Each extractor populates :class:`DocumentStructure` (headings, paragraphs,
tables, list items) and :class:`DocumentMetadata` (author, created time,
page count, etc.) when the underlying library exposes them.
"""

from __future__ import annotations

import email
import email.policy
import logging
import re
from email.utils import parsedate_to_datetime
from pathlib import Path
from typing import Any

from doctoragent._utils import truncate_text
from doctoragent.model.extensions import TEXT_EXTENSIONS, mime_type
from doctoragent.model.extractors.base import (
    MAX_EXTRACTION_CHARS,
    ContentExtractor,
    DocumentMetadata,
    DocumentStructure,
    ExtractionResult,
)

logger = logging.getLogger(__name__)

_PDF_MIME = "application/pdf"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_RTF_MIME = "application/rtf"
_EML_MIME = "message/rfc822"

_EXT_PDF = frozenset({"pdf"})
_EXT_DOCX = frozenset({"docx"})
_EXT_XLSX = frozenset({"xlsx"})
_EXT_PPTX = frozenset({"pptx"})
_EXT_EML = frozenset({"eml", "msg"})
_EXT_RTF = frozenset({"rtf"})
_EXT_PLAIN = frozenset({"txt", "md", "csv", "json", "log", "xml", "html"})

# Markdown ATX heading regex: 1-6 ``#`` followed by text.
_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$", re.MULTILINE)
# Markdown list item regex: ``-``, ``*``, ``+`` or ``N.`` prefixes.
_MD_LIST_RE = re.compile(r"^\s*(?:[-*+]|\d+\.)\s+(.*)$", re.MULTILINE)


def _import_pypdf() -> Any:
    """Lazily import pypdf; returns ``None`` when missing."""
    try:
        import pypdf  # type: ignore[import-not-found]
    except ImportError:
        return None
    return pypdf


def _import_pdfplumber() -> Any:
    """Lazily import pdfplumber; returns ``None`` when missing."""
    try:
        import pdfplumber  # type: ignore[import-not-found]
    except ImportError:
        return None
    return pdfplumber


def _import_docx() -> Any:
    """Lazily import python-docx; returns ``None`` when missing."""
    try:
        import docx  # type: ignore[import-not-found]
    except ImportError:
        return None
    return docx


def _import_openpyxl() -> Any:
    """Lazily import openpyxl; returns ``None`` when missing."""
    try:
        import openpyxl  # type: ignore[import-not-found]
    except ImportError:
        return None
    return openpyxl


def _import_pptx() -> Any:
    """Lazily import python-pptx; returns ``None`` when missing."""
    try:
        import pptx  # type: ignore[import-not-found]
    except ImportError:
        return None
    return pptx


def _import_striprtf() -> Any:
    """Lazily import striprtf; returns ``None`` when missing."""
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore[import-not-found]
    except ImportError:
        return None
    return rtf_to_text


def _word_count(text: str) -> int:
    """Return a rough word count (whitespace-split, non-empty tokens)."""
    return len([w for w in text.split() if w])


def _parse_markdown_structure(text: str) -> DocumentStructure:
    """Parse a markdown string into headings, paragraphs, and list items.

    Headings are detected via ATX syntax (``#``..``######``). List items are
    detected via ``-``/``*``/``+``/``N.`` prefixes. Paragraphs are the
    non-heading, non-list blocks of text.
    """
    structure = DocumentStructure()
    lines = text.split("\n")
    paragraph_buf: list[str] = []

    def _flush_paragraph() -> None:
        if paragraph_buf:
            para = "\n".join(paragraph_buf).strip()
            if para:
                structure.paragraphs.append(para)
            paragraph_buf.clear()

    for line in lines:
        heading_match = _MD_HEADING_RE.match(line)
        if heading_match:
            _flush_paragraph()
            level = len(heading_match.group(1))
            heading_text = heading_match.group(2).strip().rstrip("#").strip()
            if heading_text:
                structure.headings.append((level, heading_text))
            continue
        list_match = _MD_LIST_RE.match(line)
        if list_match:
            _flush_paragraph()
            item_text = list_match.group(1).strip()
            if item_text:
                structure.list_items.append(item_text)
            continue
        paragraph_buf.append(line)
    _flush_paragraph()
    return structure


def _rows_to_markdown_table(rows: list[list[Any]]) -> str:
    """Render a list of rows as a markdown table string.

    The first row is treated as the header. Empty input yields an empty
    string. Cell values are stringified; ``None`` becomes an empty cell.
    """
    if not rows:
        return ""
    str_rows = [[("" if cell is None else str(cell)) for cell in row] for row in rows]
    # Normalize ragged rows to the max column count.
    width = max(len(r) for r in str_rows)
    for r in str_rows:
        while len(r) < width:
            r.append("")
    header = str_rows[0]
    body = str_rows[1:] if len(str_rows) > 1 else []
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


class TextContentExtractor(ContentExtractor):
    """Text / document content extractor."""

    @staticmethod
    def _truncate(text: str, max_chars: int) -> tuple[str, bool]:
        return truncate_text(text, max_chars)

    def supports(self, path: Path) -> bool:
        ext = path.suffix.lower().lstrip(".")
        return ext in TEXT_EXTENSIONS or ext in _EXT_DOCX

    def extract(self, path: Path, max_chars: int = MAX_EXTRACTION_CHARS) -> ExtractionResult:
        ext = path.suffix.lower().lstrip(".")
        if ext in _EXT_PLAIN:
            return self._extract_plain(path, ext, max_chars)
        if ext in _EXT_PDF:
            return self._extract_pdf(path, max_chars)
        if ext in _EXT_DOCX:
            return self._extract_docx(path, max_chars)
        if ext in _EXT_XLSX:
            return self._extract_xlsx(path, max_chars)
        if ext in _EXT_PPTX:
            return self._extract_pptx(path, max_chars)
        if ext in _EXT_EML:
            return self._extract_eml(path, max_chars)
        if ext in _EXT_RTF:
            return self._extract_rtf(path, max_chars)
        return ExtractionResult(text="", method="none")

    def _extract_plain(self, path: Path, ext: str, max_chars: int) -> ExtractionResult:
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                text = path.read_text(encoding="latin-1")
            except OSError as exc:
                logger.debug("读取纯文本文件失败 %s: %s", path, exc)
                return ExtractionResult(text="", method="none")
        except OSError as exc:
            logger.debug("读取纯文本文件失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none")
        text, truncated = self._truncate(text, max_chars)
        structure: DocumentStructure | None = None
        if ext == "md":
            structure = _parse_markdown_structure(text)
        metadata = DocumentMetadata(word_count=_word_count(text))
        return ExtractionResult(
            text=text,
            method="text",
            mime_type=mime_type(ext),
            char_count=len(text),
            truncated=truncated,
            structure=structure,
            metadata=metadata,
        )

    # ── PDF ──────────────────────────────────────────────────────────────

    def _extract_pdf(self, path: Path, max_chars: int) -> ExtractionResult:
        """Extract PDF text with optional table/layout/metadata enrichment.

        Prefers ``pdfplumber`` for tables, page markers, and metadata; falls
        back to ``pypdf`` for plain text when pdfplumber is unavailable. When
        both are missing, degrades to ``method="none"``.
        """
        pdfplumber = _import_pdfplumber()
        if pdfplumber is not None:
            return self._extract_pdf_with_pdfplumber(path, max_chars, pdfplumber)
        return self._extract_pdf_with_pypdf(path, max_chars)

    def _extract_pdf_with_pdfplumber(
        self, path: Path, max_chars: int, pdfplumber: Any
    ) -> ExtractionResult:
        """Advanced PDF extraction: tables, page markers, metadata."""
        try:
            with pdfplumber.open(str(path)) as pdf:
                parts: list[str] = []
                tables_md: list[str] = []
                paragraphs: list[str] = []
                page_count = len(pdf.pages)
                for idx, page in enumerate(pdf.pages, start=1):
                    parts.append(f"[Page {idx}]")
                    page_text = page.extract_text() or ""
                    if page_text:
                        parts.append(page_text)
                        # Treat consecutive non-empty lines joined as paragraphs.
                        for block in page_text.split("\n\n"):
                            block = block.strip()
                            if block:
                                paragraphs.append(block)
                    # Extract tables on this page.
                    try:
                        page_tables = page.extract_tables() or []
                    except Exception:  # noqa: BLE001 — table extraction is best-effort
                        page_tables = []
                    for tbl in page_tables:
                        if tbl:
                            md = _rows_to_markdown_table(tbl)
                            if md:
                                tables_md.append(md)
                                parts.append(md)
                text = "\n".join(parts)
                text, truncated = self._truncate(text, max_chars)
                # Metadata
                meta = pdf.metadata or {}
                metadata = DocumentMetadata(
                    title=meta.get("Title") or None,
                    author=meta.get("Author") or None,
                    created_at=str(meta["CreationDate"]) if meta.get("CreationDate") else None,
                    modified_at=str(meta["ModDate"]) if meta.get("ModDate") else None,
                    page_count=page_count,
                    word_count=_word_count(text),
                    producer=meta.get("Producer") or None,
                    # Keywords are sometimes present but not in our schema; skip.
                )
                structure = DocumentStructure(
                    paragraphs=paragraphs,
                    tables=tables_md,
                )
                return ExtractionResult(
                    text=text,
                    method="text",
                    mime_type=_PDF_MIME,
                    char_count=len(text),
                    truncated=truncated,
                    structure=structure,
                    metadata=metadata,
                )
        except Exception as exc:  # noqa: BLE001
            logger.debug("pdfplumber 提取失败 %s，回退到 pypdf: %s", path, exc)
            return self._extract_pdf_with_pypdf(path, max_chars)

    def _extract_pdf_with_pypdf(self, path: Path, max_chars: int) -> ExtractionResult:
        """Fallback PDF extraction using pypdf (plain text + page markers + metadata)."""
        pypdf = _import_pypdf()
        if pypdf is None:
            logger.debug("pypdf 未安装，跳过 PDF 提取: %s", path)
            return ExtractionResult(text="", method="none", mime_type=_PDF_MIME)
        try:
            reader = pypdf.PdfReader(str(path))
            parts: list[str] = []
            paragraphs: list[str] = []
            page_count = len(reader.pages)
            for idx, page in enumerate(reader.pages, start=1):
                parts.append(f"[Page {idx}]")
                page_text = page.extract_text() or ""
                parts.append(str(page_text))
                for block in str(page_text).split("\n\n"):
                    block = block.strip()
                    if block:
                        paragraphs.append(block)
            text = "\n".join(parts)
            text, truncated = self._truncate(text, max_chars)
            # Best-effort metadata extraction.
            metadata = DocumentMetadata(page_count=page_count, word_count=_word_count(text))
            try:
                meta = reader.metadata
                if meta is not None:
                    metadata.title = getattr(meta, "title", None) or None
                    metadata.author = getattr(meta, "author", None) or None
                    metadata.producer = getattr(meta, "producer", None) or None
                    created = getattr(meta, "creation_date", None)
                    if created is not None:
                        metadata.created_at = str(created)
                    modified = getattr(meta, "modification_date", None)
                    if modified is not None:
                        metadata.modified_at = str(modified)
            except Exception:  # noqa: BLE001 — metadata is best-effort
                pass
            structure = DocumentStructure(paragraphs=paragraphs)
            return ExtractionResult(
                text=text,
                method="text",
                mime_type=_PDF_MIME,
                char_count=len(text),
                truncated=truncated,
                structure=structure,
                metadata=metadata,
            )
        except Exception as exc:  # noqa: BLE001
            logger.debug("PDF 提取失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_PDF_MIME)

    # ── DOCX ─────────────────────────────────────────────────────────────

    def _extract_docx(self, path: Path, max_chars: int) -> ExtractionResult:
        """Extract DOCX text with heading hierarchy and metadata."""
        docx = _import_docx()
        if docx is None:
            logger.debug("python-docx 未安装，跳过 DOCX 提取: %s", path)
            return ExtractionResult(text="", method="none", mime_type=_DOCX_MIME)
        try:
            document = docx.Document(str(path))
            parts: list[str] = []
            structure = DocumentStructure()
            for para in document.paragraphs:
                para_text = (para.text or "").strip()
                if not para_text:
                    continue
                parts.append(para_text)
                style_name = (para.style.name or "").lower() if para.style else ""
                # Map Word heading styles to levels.
                if style_name.startswith("heading"):
                    try:
                        level = int(style_name.replace("heading", "").strip())
                    except ValueError:
                        level = 1
                    structure.headings.append((level, para_text))
                elif style_name.startswith("list"):
                    structure.list_items.append(para_text)
                else:
                    structure.paragraphs.append(para_text)
            # Extract tables.
            for table in document.tables:
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                md = _rows_to_markdown_table(rows)
                if md:
                    structure.tables.append(md)
                    parts.append(md)
            text = "\n".join(parts)
            text, truncated = self._truncate(text, max_chars)
            # Metadata
            metadata = DocumentMetadata(word_count=_word_count(text))
            try:
                cp = document.core_properties
                metadata.author = cp.author or None
                metadata.title = cp.title or None
                metadata.created_at = cp.created.isoformat() if cp.created else None
                metadata.modified_at = cp.modified.isoformat() if cp.modified else None
                metadata.producer = cp.last_modified_by or None
            except Exception:  # noqa: BLE001 — metadata is best-effort
                pass
            return ExtractionResult(
                text=text,
                method="text",
                mime_type=_DOCX_MIME,
                char_count=len(text),
                truncated=truncated,
                structure=structure,
                metadata=metadata,
            )
        except Exception as exc:  # noqa: BLE001
            logger.debug("DOCX 提取失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_DOCX_MIME)

    # ── XLSX ─────────────────────────────────────────────────────────────

    def _extract_xlsx(self, path: Path, max_chars: int) -> ExtractionResult:
        """Extract XLSX cell text, per-sheet tables, and workbook metadata."""
        openpyxl = _import_openpyxl()
        if openpyxl is None:
            logger.debug("openpyxl 未安装，跳过 XLSX 提取: %s", path)
            return ExtractionResult(text="", method="none", mime_type=_XLSX_MIME)
        try:
            wb = openpyxl.load_workbook(filename=str(path), read_only=True, data_only=True)
            parts: list[str] = []
            tables_md: list[str] = []
            sheet_names = wb.sheetnames
            for name in sheet_names:
                ws = wb[name]
                parts.append(f"[Sheet: {name}]")
                rows: list[list[Any]] = []
                for row in ws.iter_rows(values_only=True):
                    rows.append(list(row))
                if rows:
                    md = _rows_to_markdown_table(rows)
                    if md:
                        tables_md.append(md)
                        parts.append(md)
            wb.close()
            text = "\n".join(parts)
            text, truncated = self._truncate(text, max_chars)
            metadata = DocumentMetadata(
                page_count=len(sheet_names),
                word_count=_word_count(text),
            )
            # Best-effort workbook properties.
            try:
                props = wb.properties
                metadata.title = props.title or None
                metadata.author = props.creator or None
                metadata.created_at = props.created.isoformat() if props.created else None
                metadata.modified_at = props.modified.isoformat() if props.modified else None
            except Exception:  # noqa: BLE001 — properties access after close is best-effort
                pass
            structure = DocumentStructure(tables=tables_md)
            return ExtractionResult(
                text=text,
                method="text",
                mime_type=_XLSX_MIME,
                char_count=len(text),
                truncated=truncated,
                structure=structure,
                metadata=metadata,
            )
        except Exception as exc:  # noqa: BLE001
            logger.debug("XLSX 提取失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_XLSX_MIME)

    # ── PPTX ─────────────────────────────────────────────────────────────

    def _extract_pptx(self, path: Path, max_chars: int) -> ExtractionResult:
        """Extract PPTX slide text, notes, and presentation metadata."""
        pptx = _import_pptx()
        if pptx is None:
            logger.debug("python-pptx 未安装，跳过 PPTX 提取: %s", path)
            return ExtractionResult(text="", method="none", mime_type=_PPTX_MIME)
        try:
            prs = pptx.Presentation(str(path))
            parts: list[str] = []
            headings: list[tuple[int, str]] = []
            paragraphs: list[str] = []
            slide_count = len(prs.slides)
            for idx, slide in enumerate(prs.slides, start=1):
                parts.append(f"[Slide {idx}]")
                title_text = ""
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        para_text = (para.text or "").strip()
                        if not para_text:
                            continue
                        parts.append(para_text)
                        # The first text shape on a slide is usually the title.
                        if not title_text and shape == slide.shapes.title:
                            title_text = para_text
                            headings.append((1, para_text))
                        else:
                            paragraphs.append(para_text)
                # Speaker notes.
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text or ""
                    if notes_text.strip():
                        parts.append(f"[Notes] {notes_text.strip()}")
            text = "\n".join(parts)
            text, truncated = self._truncate(text, max_chars)
            metadata = DocumentMetadata(
                page_count=slide_count,
                word_count=_word_count(text),
            )
            # Best-effort core properties.
            try:
                cp = prs.core_properties
                metadata.title = cp.title or None
                metadata.author = cp.author or None
                metadata.created_at = cp.created.isoformat() if cp.created else None
                metadata.modified_at = cp.modified.isoformat() if cp.modified else None
            except Exception:  # noqa: BLE001 — metadata is best-effort
                pass
            structure = DocumentStructure(headings=headings, paragraphs=paragraphs)
            return ExtractionResult(
                text=text,
                method="text",
                mime_type=_PPTX_MIME,
                char_count=len(text),
                truncated=truncated,
                structure=structure,
                metadata=metadata,
            )
        except Exception as exc:  # noqa: BLE001
            logger.debug("PPTX 提取失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_PPTX_MIME)

    # ── EML ──────────────────────────────────────────────────────────────

    def _extract_eml(self, path: Path, max_chars: int) -> ExtractionResult:
        """Extract EML email headers, body, and attachment names.

        Uses the stdlib ``email`` module so no third-party dependency is
        required. Falls back to ``method="none"`` only on parse failure.
        """
        try:
            raw = path.read_bytes()
        except OSError as exc:
            logger.debug("读取 EML 文件失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_EML_MIME)
        try:
            msg = email.message_from_bytes(raw, policy=email.policy.default)
        except Exception as exc:  # noqa: BLE001
            logger.debug("EML 解析失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_EML_MIME)
        try:
            subject = str(msg.get("Subject", "") or "").strip()
            from_ = str(msg.get("From", "") or "").strip()
            to_ = str(msg.get("To", "") or "").strip()
            date_raw = msg.get("Date", "")
            parts: list[str] = []
            if subject:
                parts.append(f"Subject: {subject}")
            if from_:
                parts.append(f"From: {from_}")
            if to_:
                parts.append(f"To: {to_}")
            if date_raw:
                parts.append(f"Date: {date_raw}")
            # Extract body (prefer plain text part).
            body_text = ""
            attachments: list[str] = []
            if msg.is_multipart():
                for part in msg.walk():
                    content_disposition = str(part.get_content_disposition() or "")
                    if content_disposition == "attachment":
                        fn = part.get_filename()
                        if fn:
                            attachments.append(fn)
                        continue
                    ctype = part.get_content_type()
                    if ctype == "text/plain":
                        try:
                            payload = part.get_content()
                        except Exception:  # noqa: BLE001
                            payload = part.get_payload(decode=True)
                            if isinstance(payload, bytes):
                                payload = payload.decode("utf-8", errors="replace")
                        body_text = str(payload or "")
                        break
                if not body_text:
                    # Fall back to first text part of any kind.
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            try:
                                payload = part.get_content()
                            except Exception:  # noqa: BLE001
                                payload = part.get_payload(decode=True)
                                if isinstance(payload, bytes):
                                    payload = payload.decode("utf-8", errors="replace")
                            body_text = str(payload or "")
                            break
            else:
                try:
                    body_text = str(msg.get_content() or "")
                except Exception:  # noqa: BLE001
                    payload = msg.get_payload(decode=True)
                    if isinstance(payload, bytes):
                        body_text = payload.decode("utf-8", errors="replace")
                    else:
                        body_text = str(payload or "")
            if body_text:
                parts.append("")
                parts.append(body_text)
            if attachments:
                parts.append("")
                parts.append("Attachments: " + ", ".join(attachments))
            text = "\n".join(parts)
            text, truncated = self._truncate(text, max_chars)
            # Metadata
            created_at: str | None = None
            if date_raw:
                try:
                    dt = parsedate_to_datetime(date_raw)
                    if dt is not None:
                        created_at = dt.isoformat()
                except (TypeError, ValueError):
                    created_at = str(date_raw)
            metadata = DocumentMetadata(
                title=subject or None,
                author=from_ or None,
                created_at=created_at,
                word_count=_word_count(text),
            )
            structure = DocumentStructure(paragraphs=[body_text] if body_text else [])
            return ExtractionResult(
                text=text,
                method="text",
                mime_type=_EML_MIME,
                char_count=len(text),
                truncated=truncated,
                structure=structure,
                metadata=metadata,
            )
        except Exception as exc:  # noqa: BLE001
            logger.debug("EML 内容提取失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_EML_MIME)

    # ── RTF ──────────────────────────────────────────────────────────────

    def _extract_rtf(self, path: Path, max_chars: int) -> ExtractionResult:
        """Extract RTF plain text via ``striprtf``."""
        rtf_to_text = _import_striprtf()
        if rtf_to_text is None:
            logger.debug("striprtf 未安装，跳过 RTF 提取: %s", path)
            return ExtractionResult(text="", method="none", mime_type=_RTF_MIME)
        try:
            raw = path.read_text(encoding="utf-8", errors="replace")
            text = rtf_to_text(raw)
            text, truncated = self._truncate(text, max_chars)
            metadata = DocumentMetadata(word_count=_word_count(text))
            structure = _parse_markdown_structure(text)
            return ExtractionResult(
                text=text,
                method="text",
                mime_type=_RTF_MIME,
                char_count=len(text),
                truncated=truncated,
                structure=structure,
                metadata=metadata,
            )
        except Exception as exc:  # noqa: BLE001
            logger.debug("RTF 提取失败 %s: %s", path, exc)
            return ExtractionResult(text="", method="none", mime_type=_RTF_MIME)
